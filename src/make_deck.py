#!/usr/bin/env python3
"""
Build the Stage-1 (mid-presentation) deck as a native .pptx.

Hebrew is right-to-left, set per paragraph via the OOXML attribute a:pPr rtl="1"
(python-pptx has no RTL property). Figures from ../figures; deck to ../presentation.
Open in Google Slides / PowerPoint / LibreOffice. Hebrew strings = slide content.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DIR = str(ROOT / "figures")
FONT = "Arial"
NAVY, BLUE, INK, GREY = RGBColor(0x1F,0x38,0x64), RGBColor(0x2E,0x86,0xC1), RGBColor(0x22,0x22,0x22), RGBColor(0x70,0x70,0x70)
HL = RGBColor(0xC0,0x53,0x1B)
SW, SH = 13.333, 7.5
LEFT, BODYW = 0.6, 12.13

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
BLANK = prs.slide_layouts[6]

def _set(p, align="r", rtl=True):
    pPr = p._p.get_or_add_pPr()
    if rtl: pPr.set("rtl", "1")
    pPr.set("algn", {"r":"r","l":"l","ctr":"ctr"}.get(align,"r"))

def textbox(slide, l, t, w, h, anchor=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor: tf.vertical_anchor = anchor
    return tf

def line(tf, text, size, color=INK, bold=False, align="r", space_after=8, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
    r.font.color.rgb = color
    _set(p, align)
    return p

def rule(slide, top, l=LEFT, w=BODYW, color=BLUE, h=0.045):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background(); s.shadow.inherit = False

def footer(slide):
    n = len(list(prs.slides))
    line(textbox(slide, LEFT, 7.08, BODYW, 0.3), f"תאוריה סטטיסטית · פרויקט סוף · OpenPowerlifting · {n}/9",
         9, GREY, align="r", space_after=0, first=True)

def header(slide, title):
    line(textbox(slide, LEFT, 0.45, BODYW, 0.95), title, 30, NAVY, bold=True, first=True)
    rule(slide, 1.45); footer(slide)

def content(title, bullets, lead=None, lead_size=22):
    s = prs.slides.add_slide(BLANK); header(s, title)
    tf = textbox(s, LEFT, 1.8, BODYW, 5.0, anchor=MSO_ANCHOR.TOP); first = True
    if lead:
        line(tf, lead, lead_size, BLUE, bold=True, space_after=18, first=True); first = False
    for b in bullets:
        sub = b.startswith("\t")
        line(tf, ("◦  " if sub else "•  ") + b.lstrip("\t"), 18 if sub else 19.5, INK, space_after=14, first=first)
        first = False
    return s

def fig_slide(title, lead, img, aspect, max_w, sub=None):
    s = prs.slides.add_slide(BLANK); header(s, title)
    y = 1.6
    if sub:
        line(textbox(s, LEFT, y, BODYW, 0.55), sub, 21, HL, bold=True, align="ctr", space_after=0, first=True)
        y = 2.18
    line(textbox(s, LEFT, y, BODYW, 1.15), lead, 18, INK, space_after=0, first=True)
    top = 2.78 if not sub else 3.30; avail_h = 6.95 - top
    w = min(max_w, avail_h*aspect); h = w/aspect
    s.shapes.add_picture(f"{DIR}/{img}", Inches((SW-w)/2), Inches(top), width=Inches(w), height=Inches(h))
    return s

# ---- Slide 1: title (cover identity + hook number) ----
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(0.55))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background(); band.shadow.inherit = False
tf = textbox(s, 1.0, 1.9, 11.33, 4.0, anchor=MSO_ANCHOR.TOP)
line(tf, 'איך מתחרי הרמת-כוח "משחקים" עם המספרים', 38, NAVY, bold=True, align="ctr", space_after=10, first=True)
line(tf, "ניתוח סטטיסטי של OpenPowerlifting: מה הנתונים חושפים על גבולות הכוח", 19, GREY, align="ctr", space_after=22)
line(tf, 'פי כ-6.8 יותר מתחרים נשקלים ממש מתחת לסף מחלקת-המשקל מאשר ממש מעליו', 21, HL, bold=True, align="ctr", space_after=26)
line(tf, "אדיר אלמקייס · דוד לוין", 22, INK, bold=True, align="ctr", space_after=6)
line(tf, "תאוריה סטטיסטית · פרויקט סוף · מצגת אמצע", 16, GREY, align="ctr")
rule(s, 6.05, l=4.4, w=4.5)

# ---- Slide 2: research data ----
content("נתוני המחקר", [
    "מקור: OpenPowerlifting, מאגר פתוח של תוצאות תחרויות הרמת-כוח (רישיון CC0).",
    "היקף: כ-3.94 מיליון תוצאות, 42 משתנים. יחידת-התצפית: שורה לכל מתחרה בכל תחרות.",
    "אתגר עיקרי (אי-תלות): אותו מתחרה מופיע בכמה תחרויות, אז במבחנים הפורמליים ניקח שורה אחת לכל מתחרה.",
    "שומרים גרסת-נתונים קבועה לשחזור.",
])

# ---- Slide 3: overview (two-numbers tagline + contribution) ----
content("עיקרי העבודה", [
    "רקע: בהרמת-כוח שלוש הרמות (סקוואט, בנץ', דדליפט), שלושה ניסיונות לכל הרמה, והתחרות בתוך קטגוריות משקל-גוף.",
    "מתחרה שולט בדיוק בשני מספרים: המשקל על המוט (בחירת ניסיונות) ומשקל הגוף (בשקילה).",
    "התרומה שלנו: לא רק להראות את התופעה, אלא לבדוק אותה בצורה מסודרת, עם מודל ועם נקודת-בקרה.",
], lead="הרעיון: שני המספרים שהמתחרה שולט בהם, המוט והמאזניים")

# ---- Slide 4: research question + hypotheses ----
content("שאלת החקר", [
    'ארבע ההשערות מפרקות את השאלה: H1/H2 הם "שני המספרים", H3/H4 הם גבולות-הכוח והניבוי.',
    "\tH1: משקלי הניסיון אינם רציפים אלא נופלים כמעט תמיד על רשת 2.5 ק\"ג.",
    "\tH2: הצטברות משקל-גוף ממש מתחת לסף קטגוריית-המשקל (עקבי עם חיתוך-משקל).",
    "\tH3: קנה-המידה האלומטרי (כוח מול משקל-גוף) שונה בין המינים.",
    "\tH4 (לעבודה הסופית): כמה מהכוח ניתן לנבא ממשתנים בסיסיים (משקל-גוף, ציוד, מין, גיל)?",
], lead='שאלת המחקר: כיצד מתחרי הרמת-כוח "משחקים" עם שני המספרים שבשליטתם (המוט והמאזניים), ומה הדפוסים מגלים על גבולות הכוח?')

# ---- Slide 5: methods (1:1 with results + named course tools + ML) ----
content("שיטות החקר", [
    "H1 (רשת המוט): מודל שבודק אם המשקלים יושבים על רשת 2.5 ק\"ג, עם מבחן GLRT (השוואת שני מודלים).",
    "H2 (חיתוך-משקל): מבחן McCrary לקפיצה בצפיפות סביב הסף, עם נקודת-בקרה וניקוי מספרים-עגולים.",
    "H3 (כוח מול משקל): רגרסיה שבודקת אם קנה-המידה שונה מ-2/3, ואם הוא שונה בין המינים (מבחן Wald).",
    "H4 (לעבודה הסופית): מודל ניבוי-כוח, רגרסיה + Random Forest + סיווג, שנבדק על נתונים חדשים.",
    'כי המדגם ענק (כמעט 4 מיליון, עוצמה גבוהה), כמעט הכל יוצא "מובהק", אז מדגישים גודל-אפקט ורווח-סמך, לא p. מיישמים גם תיקון להשוואות מרובות.',
], lead_size=20)

# ---- Slide 6: results H1 ----
fig_slide('תוצאות ראשוניות (1): רשת ה-2.5 ק"ג',
          'משקלי הניסיון אינם רציפים: 96.2% מתוך 13.4 מיליון ניסיונות בדיוק על רשת ה-2.5 ק"ג ‎(95% CI [96.19%, 96.21%])‎.',
          "fig_quantization.png", 8.2/4.4, 9.4)

# ---- Slide 7: results H2 ----
fig_slide("תוצאות ראשוניות (2): חיתוך-משקל",
          'פי כ-6.8 יותר מתחת לסף 83 ק"ג מאשר מעליו (גברים, IPF+USAPL): יחס-לוג ‎+1.92 ± 0.04‎, גם אחרי ניקוי מספרים-עגולים. בבקרה 91 ק"ג: ‎−0.21 ± 0.03‎. עקבי עם חיתוך-משקל; אישוש פורמלי (McCrary) בהמשך.',
          "fig_bunching.png", 9.8/4.3, 11.2,
          sub='סף אמיתי 83 ק"ג → הצטברות · בקרה 91 ק"ג → אין')

# ---- Slide 8: results H3 (allometry) ----
fig_slide("תוצאות ראשוניות (3): קנה-מידה אלומטרי",
          'הכוח גדל עם משקל-הגוף, אבל אחרת אצל גברים ונשים: ‎b≈0.72 (R²=0.36)‎ גברים / ‎0.49 (R²=0.19)‎ נשים, מול 2/3 בתאוריה. זה אומדן ראשוני; נבדוק פורמלית במבחן Wald.',
          "fig_allometry.png", 9.8/4.3, 11.2)

# ---- Slide 9: conclusions + payoff close ----
s = prs.slides.add_slide(BLANK); header(s, "מסקנות ראשוניות וסיכום")
tf = textbox(s, LEFT, 1.75, BODYW, 4.6)
bullets = [
    'שני המספרים שהמתחרה שולט בהם נראים בבירור: רשת ה-2.5 ק"ג (המוט), והצטברות מתחת לספים (משקל גוף).',
    "ההצטברות שורדת ניקוי-עיגול ונעדרת בנקודת-בקרה: אינדיקציה ראשונית מעודדת (טרם מבוססת).",
    "הצעדים הבאים: McCrary פורמלי + בדיקות-עמידות (פלצבו), אלומטריה (Wald), ומודל ניבוי-כוח (רגרסיה + Random Forest, R²/CV).",
]
for i, b in enumerate(bullets):
    line(tf, "•  " + b, 19, INK, space_after=13, first=(i == 0))
line(tf, "המסקנה: האסטרטגיה מותירה עקבות מדידות בנתונים, ואנחנו עוברים מתיאור, למבחן, לניבוי.", 19, NAVY, bold=True, space_after=14)
line(tf, "תודה! שאלות?", 24, BLUE, bold=True, align="ctr", space_after=0)

# ---- Slide 10: BACKUP (not shown in the 5-min flow; only if asked about assumptions) ----
s = prs.slides.add_slide(BLANK)
line(textbox(s, LEFT, 0.45, BODYW, 0.95), 'גיבוי · למה לא מבחני-נורמליות פורמליים', 28, NAVY, bold=True, first=True)
rule(s, 1.45)
line(textbox(s, LEFT, 7.08, BODYW, 0.3),
     "שקף גיבוי · לא מוצג במצגת הראשית · נשלף רק לשאלה על הנחות-המודל", 9, GREY, align="r", space_after=0, first=True)
line(textbox(s, LEFT, 1.6, BODYW, 1.35),
     'המבחנים המרכזיים שלנו לא נשענים על נורמליות: GLRT על מודל בדיד (H1), McCrary על צפיפות (H2), '
     'ולמידת-מכונה (H4). ברגרסיית האלומטריה (H3) השאריות אינן נורמליות לחלוטין (זנב-תחתון כבד), '
     'אבל אמידת-השיפוע תקפה בזכות משפט-הגבול-המרכזי, ומבחן-נורמליות פורמלי דוחה ב-n ענק כל סטייה זעירה '
     '(בדיוק כמו p): לכן מובילים בגודל-אפקט וברווח-סמך.', 16, INK, space_after=0, first=True)
_img, _asp, _maxw = "fig_normality_backup.png", 10.0/4.3, 11.8
_top = 3.25; _availh = 6.85 - _top
_w = min(_maxw, _availh*_asp); _h = _w/_asp
s.shapes.add_picture(f"{DIR}/{_img}", Inches((SW-_w)/2), Inches(_top), width=Inches(_w), height=Inches(_h))

out = str(ROOT / "presentation" / "mid_presentation_OpenPowerlifting.pptx")
prs.save(out)
print("saved:", out, "|", len(list(prs.slides)), "slides")
